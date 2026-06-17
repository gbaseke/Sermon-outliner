#!/usr/bin/env python
"""Build a sermon PowerPoint from a JSON slide spec, using the bundled template.

Usage:
    uv run python build_pptx.py <spec.json>

The spec is a JSON object:
{
  "output": "C:/.../output/my-sermon.pptx",   # required, where to write the .pptx
  "template": "C:/.../template.pptx",          # optional, defaults to template.pptx next to this script
  "slides": [ <slide>, ... ]                   # ordered list of slide objects
}

Slide objects (the "type" field selects the template layout):

  {"type": "title",     "title": str, "subtitle": str?, "footer": str?, "date": str?}
  {"type": "scripture", "reference": str, "text": str}        # text may hold several verses
  {"type": "quote",     "text": str, "citation": str?}        # Branham / message quote
  {"type": "notes",     "header": str, "bullets": [str, ...]} # bulleted note card
  {"type": "point",     "title": str, "bullets": [str, ...]?} # Title + Content (sub-points)
  {"type": "section",   "title": str}                         # point/section divider
  {"type": "blank"}                                           # empty slide

Pagination is computed from the *real* placeholder geometry and font size in the
template (body text is 40pt), so long scripture / quote / notes / point slides are
split across as many slides as it takes to fit -- nothing overflows the box.
Split slides get a "(n/m)" marker.
"""

import json
import math
import re
import sys
from pathlib import Path

from pptx import Presentation
from pptx.oxml.ns import qn

# --- layout names as they appear in template.pptx ---
L_TITLE = "Title Slide"
L_QUOTE = "Quote"
L_SCRIPTURE = "Scripture"
L_NOTES = "Notes"
L_CONTENT = "Title and Content"
L_BLANK = "Blank"

# placeholder indices within each layout
PH_TITLE_MAIN, PH_TITLE_SUB = 0, 1        # Title Slide
PH_HEADER, PH_BODY = 13, 14               # Quote / Scripture / Notes
PH_CONTENT_TITLE, PH_CONTENT_BODY = 0, 1  # Title and Content

EMU_PER_IN = 914400

# How wide an average glyph is, and how tall a line is, as a fraction of the
# font's point size. Tuned conservatively for the template's serif body font so
# we under-fill rather than overflow.
CHAR_W_FACTOR = 0.52
LINE_H_FACTOR = 1.25
FILL_SAFETY = 0.90       # never trust the last 10% of a full line
DEFAULT_BODY_PT = 40     # fallback if the template font size can't be read

# A "section" title longer than this many characters won't fit the 1-line title
# box, so it is rendered as a big centered statement on the Quote layout instead.
SECTION_TITLE_MAX_CHARS = 38


def layout_by_name(prs, name):
    for lyt in prs.slide_layouts:
        if lyt.name == name:
            return lyt
    raise SystemExit(f"Template is missing layout '{name}'. "
                     f"Available: {[l.name for l in prs.slide_layouts]}")


def ph_by_idx(container, idx):
    for ph in container.placeholders:
        if ph.placeholder_format.idx == idx:
            return ph
    return None


def body_font_pt(prs):
    """Read the master body style's level-1 font size (the body text size)."""
    try:
        master = prs.slide_masters[0].element
        txstyles = master.find(qn("p:txStyles"))
        lvl1 = txstyles.find(qn("p:bodyStyle")).find(qn("a:lvl1pPr"))
        sz = lvl1.find(qn("a:defRPr")).get("sz")
        return int(sz) / 100
    except Exception:
        return DEFAULT_BODY_PT


def capacity(width_emu, height_emu, font_pt):
    """Return (chars_per_line, max_lines) that fit a box at the given font size."""
    w_in = width_emu / EMU_PER_IN
    h_in = height_emu / EMU_PER_IN
    cpl = max(8, int(w_in / (font_pt * CHAR_W_FACTOR / 72)))
    lines = max(1, int(h_in / (font_pt * LINE_H_FACTOR / 72)))
    return cpl, lines


def body_capacity(layout, idx, font_pt):
    ph = ph_by_idx(layout, idx)
    return capacity(ph.width, ph.height, font_pt)


def set_notes(slide, text):
    """Put preacher prose in the slide's speaker-notes pane (off the screen)."""
    if not text:
        return
    slide.notes_slide.notes_text_frame.text = str(text)


def set_text(placeholder, text):
    """Set a single block of text, preserving the placeholder's template style."""
    if placeholder is None:
        return
    placeholder.text_frame.word_wrap = True
    placeholder.text_frame.text = "" if text is None else str(text)


def set_paragraphs(placeholder, lines):
    """Fill a placeholder with multiple paragraphs, inheriting template formatting."""
    if placeholder is None:
        return
    tf = placeholder.text_frame
    tf.word_wrap = True
    lines = [l for l in lines if l is not None]
    if not lines:
        tf.text = ""
        return
    tf.text = str(lines[0])
    for line in lines[1:]:
        p = tf.add_paragraph()
        p.text = str(line)


# ---------------------------------------------------------------------------
# text-splitting helpers (line-aware, driven by real box capacity)
# ---------------------------------------------------------------------------

def wrapped_lines(text, cpl):
    """Estimate how many wrapped lines a single paragraph occupies."""
    text = (text or "").strip()
    return max(1, math.ceil(len(text) / cpl)) if text else 1


def split_text_to_chunks(text, max_chars):
    """Split a long string into <=max_chars chunks, preferring sentence breaks."""
    text = (text or "").strip()
    if len(text) <= max_chars:
        return [text] if text else [""]
    sentences = re.split(r"(?<=[.!?;:])\s+", text)
    chunks, cur = [], ""
    for s in sentences:
        if len(s) > max_chars:                   # a single huge sentence -> split by words
            if cur:
                chunks.append(cur.strip())
                cur = ""
            line = ""
            for w in s.split():
                if len(line) + len(w) + 1 > max_chars:
                    chunks.append(line.strip())
                    line = w
                else:
                    line = f"{line} {w}".strip()
            if line:
                cur = line
            continue
        if len(cur) + len(s) + 1 > max_chars:
            chunks.append(cur.strip())
            cur = s
        else:
            cur = f"{cur} {s}".strip()
    if cur.strip():
        chunks.append(cur.strip())
    return chunks


def split_block(text, cpl, max_lines):
    """Split a prose block so each chunk fits within max_lines at cpl chars/line."""
    budget = int(cpl * max_lines * FILL_SAFETY)
    return split_text_to_chunks(text, budget)


def split_paragraphs(paragraphs, cpl, max_lines):
    """Group paragraphs (e.g. verses) so each slide's wrapped lines <= max_lines."""
    chunks, cur, cur_lines = [], [], 0
    for para in paragraphs:
        para = para.strip()
        if not para:
            continue
        pl = wrapped_lines(para, cpl)
        if pl > max_lines:                       # one paragraph too big on its own
            if cur:
                chunks.append(cur)
                cur, cur_lines = [], 0
            for piece in split_block(para, cpl, max_lines):
                chunks.append([piece])
            continue
        if cur and cur_lines + pl > max_lines:
            chunks.append(cur)
            cur, cur_lines = [], 0
        cur.append(para)
        cur_lines += pl
    if cur:
        chunks.append(cur)
    return chunks or [[]]


def split_bullets(bullets, cpl, max_lines):
    """Group bullets so each slide's total wrapped lines fits (reserving spacing)."""
    usable = max(2, max_lines - 1)               # reserve a line for bullet spacing
    chunks, cur, cur_lines = [], [], 0
    for b in bullets:
        b = (b or "").strip()
        if not b:
            continue
        bl = wrapped_lines(b, cpl)
        if cur and cur_lines + bl > usable:
            chunks.append(cur)
            cur, cur_lines = [], 0
        cur.append(b)
        cur_lines += bl
    if cur:
        chunks.append(cur)
    return chunks or [[]]


def marker(title, i, total):
    """Append a (n/m) marker when content was split across multiple slides."""
    if not title:
        return ""
    return title if total <= 1 else f"{title}  ({i + 1}/{total})"


# ---------------------------------------------------------------------------
# per-type slide builders
# ---------------------------------------------------------------------------

def add_title(prs, layouts, caps, s):
    slide = prs.slides.add_slide(layouts[L_TITLE])
    set_text(ph_by_idx(slide, PH_TITLE_MAIN), s.get("title", ""))
    # The title layout renders only Title + Subtitle (date/footer placeholders are
    # master-driven and not cloned), so author lines fold into the subtitle.
    sub_lines = [s.get("subtitle"), s.get("footer"), s.get("date")]
    set_paragraphs(ph_by_idx(slide, PH_TITLE_SUB), [l for l in sub_lines if l])


def add_scripture(prs, layouts, caps, s):
    cpl, max_lines = caps[L_SCRIPTURE]
    ref = s.get("reference", "")
    paras = re.split(r"\n+", (s.get("text") or "").strip())
    chunks = split_paragraphs(paras, cpl, max_lines)
    for i, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(layouts[L_SCRIPTURE])
        set_text(ph_by_idx(slide, PH_HEADER), marker(ref, i, len(chunks)))
        set_paragraphs(ph_by_idx(slide, PH_BODY), chunk)


def add_quote(prs, layouts, caps, s):
    cpl, max_lines = caps[L_QUOTE]
    citation = s.get("citation", "")
    chunks = split_block(s.get("text", ""), cpl, max_lines)
    for i, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(layouts[L_QUOTE])
        set_text(ph_by_idx(slide, PH_HEADER), marker(citation, i, len(chunks)))
        set_text(ph_by_idx(slide, PH_BODY), chunk)


def add_notes(prs, layouts, caps, s):
    cpl, max_lines = caps[L_NOTES]
    header = s.get("header", "")
    chunks = split_bullets(s.get("bullets", []), cpl, max_lines)
    for i, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(layouts[L_NOTES])
        set_text(ph_by_idx(slide, PH_HEADER), marker(header, i, len(chunks)))
        set_paragraphs(ph_by_idx(slide, PH_BODY), chunk)


def add_point(prs, layouts, caps, s):
    cpl, max_lines = caps[L_CONTENT]
    title = s.get("title", "")
    bullets = s.get("bullets", [])
    if not bullets:
        slide = prs.slides.add_slide(layouts[L_CONTENT])
        set_text(ph_by_idx(slide, PH_CONTENT_TITLE), title)
        return
    chunks = split_bullets(bullets, cpl, max_lines)
    for i, chunk in enumerate(chunks):
        slide = prs.slides.add_slide(layouts[L_CONTENT])
        set_text(ph_by_idx(slide, PH_CONTENT_TITLE), marker(title, i, len(chunks)))
        set_paragraphs(ph_by_idx(slide, PH_CONTENT_BODY), chunk)


def add_section(prs, layouts, caps, s):
    title = s.get("title", "")
    # A short divider fits the 1-line title box; a long statement (e.g. the
    # "heart of the message") is rendered big on the Quote body so it never clips.
    if len(title) <= SECTION_TITLE_MAX_CHARS:
        slide = prs.slides.add_slide(layouts[L_CONTENT])
        set_text(ph_by_idx(slide, PH_CONTENT_TITLE), title)
    else:
        cpl, max_lines = caps[L_QUOTE]
        chunks = split_block(title, cpl, max_lines)
        for chunk in chunks:
            slide = prs.slides.add_slide(layouts[L_QUOTE])
            set_text(ph_by_idx(slide, PH_BODY), chunk)


def add_blank(prs, layouts, caps, s):
    prs.slides.add_slide(layouts[L_BLANK])


BUILDERS = {
    "title": add_title,
    "scripture": add_scripture,
    "quote": add_quote,
    "notes": add_notes,
    "point": add_point,
    "section": add_section,
    "blank": add_blank,
}


def main():
    if len(sys.argv) != 2:
        raise SystemExit("Usage: python build_pptx.py <spec.json>")

    spec = json.loads(Path(sys.argv[1]).read_text(encoding="utf-8"))

    template = spec.get("template")
    template_path = Path(template) if template else Path(__file__).with_name("template.pptx")
    if not template_path.exists():
        raise SystemExit(f"Template not found: {template_path}")

    out_path = Path(spec["output"])
    out_path.parent.mkdir(parents=True, exist_ok=True)

    prs = Presentation(str(template_path))
    layouts = {name: layout_by_name(prs, name) for name in
               (L_TITLE, L_QUOTE, L_SCRIPTURE, L_NOTES, L_CONTENT, L_BLANK)}

    font_pt = body_font_pt(prs)
    caps = {
        L_QUOTE: body_capacity(layouts[L_QUOTE], PH_BODY, font_pt),
        L_SCRIPTURE: body_capacity(layouts[L_SCRIPTURE], PH_BODY, font_pt),
        L_NOTES: body_capacity(layouts[L_NOTES], PH_BODY, font_pt),
        L_CONTENT: body_capacity(layouts[L_CONTENT], PH_CONTENT_BODY, font_pt),
    }

    for n, s in enumerate(spec.get("slides", [])):
        t = s.get("type")
        if t not in BUILDERS:
            raise SystemExit(f"Slide {n}: unknown type '{t}'. "
                             f"Valid types: {sorted(BUILDERS)}")
        before = len(prs.slides._sldIdLst)
        BUILDERS[t](prs, layouts, caps, s)
        # Optional "notes": preacher prose for the speaker-notes pane, attached to
        # the first slide this entry produced (the screen text stays clean).
        if s.get("notes"):
            set_notes(prs.slides[before], s["notes"])

    prs.save(str(out_path))
    print(f"Body font {font_pt:g}pt -> capacities {caps}")
    print(f"Wrote {len(list(prs.slides))} slides -> {out_path}")


if __name__ == "__main__":
    main()

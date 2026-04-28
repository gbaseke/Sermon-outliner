from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

passages = [
    ("Exodus 16:14-35", [
        (14, "And when the dew that lay was gone up, behold, upon the face of the wilderness there lay a small round thing, as small as the hoar frost on the ground."),
        (15, "And when the children of Israel saw it, they said one to another, It is manna: for they wist not what it was. And Moses said unto them, This is the bread which the Lord hath given you to eat."),
        (16, "This is the thing which the Lord hath commanded, Gather of it every man according to his eating, an omer for every man, according to the number of your persons; take ye every man for them which are in his tents."),
        (17, "And the children of Israel did so, and gathered, some more, some less."),
        (18, "And when they did mete it with an omer, he that gathered much had nothing over, and he that gathered little had no lack; they gathered every man according to his eating."),
        (19, "And Moses said, Let no man leave of it till the morning."),
        (20, "Notwithstanding they hearkened not unto Moses; but some of them left of it until the morning, and it bred worms, and stank: and Moses was wroth with them."),
        (21, "And they gathered it every morning, every man according to his eating: and when the sun waxed hot, it melted."),
        (22, "And it came to pass, that on the sixth day they gathered twice as much bread, two omers for one man: and all the rulers of the congregation came and told Moses."),
        (23, "And he said unto them, This is that which the Lord hath said, To morrow is the rest of the holy sabbath unto the Lord: bake that which ye will bake to day, and seethe that ye will seethe; and that which remaineth over lay up for you to be kept until the morning."),
        (24, "And they laid it up till the morning, as Moses bade: and it did not stink, neither was there any worm therein."),
        (25, "And Moses said, Eat that to day; for to day is a sabbath unto the Lord: to day ye shall not find it in the field."),
        (26, "Six days ye shall gather it; but on the seventh day, which is the sabbath, in it there shall be none."),
        (27, "And it came to pass, that there went out some of the people on the seventh day for to gather, and they found none."),
        (28, "And the Lord said unto Moses, How long refuse ye to keep my commandments and my laws?"),
        (29, "See, for that the Lord hath given you the sabbath, therefore he giveth you on the sixth day the bread of two days; abide ye every man in his place, let no man go out of his place on the seventh day."),
        (30, "So the people rested on the seventh day."),
        (31, "And the house of Israel called the name thereof Manna: and it was like coriander seed, white; and the taste of it was like wafers made with honey."),
        (32, "And Moses said, This is the thing which the Lord commandeth, Fill an omer of it to be kept for your generations; that they may see the bread wherewith I have fed you in the wilderness, when I brought you forth from the land of Egypt."),
        (33, "And Moses said unto Aaron, Take a pot, and put an omer full of manna therein, and lay it up before the Lord, to be kept for your generations."),
        (34, "As the Lord commanded Moses, so Aaron laid it up before the Testimony, to be kept."),
        (35, "And the children of Israel did eat manna forty years, until they came to a land inhabited; they did eat manna, until they came unto the borders of the land of Canaan."),
    ]),
    ("Numbers 11:6-9", [
        (6, "But now our soul is dried away: there is nothing at all, beside this manna, before our eyes."),
        (7, "And the manna was as coriander seed, and the colour thereof as the colour of bdellium."),
        (8, "And the people went about, and gathered it, and ground it in mills, or beat it in a mortar, and baked it in pans, and made cakes of it: and the taste of it was as the taste of fresh oil."),
        (9, "And when the dew fell upon the camp in the night, the manna fell upon it."),
    ]),
    ("John 6:26-71", [
        (26, "Jesus answered them and said, Verily, verily, I say unto you, Ye seek me, not because ye saw the miracles, but because ye did eat of the loaves, and were filled."),
        (27, "Labour not for the meat which perisheth, but for that meat which endureth unto everlasting life, which the Son of man shall give unto you: for him hath God the Father sealed."),
        (28, "Then said they unto him, What shall we do, that we might work the works of God?"),
        (29, "Jesus answered and said unto them, This is the work of God, that ye believe on him whom he hath sent."),
        (30, "They said therefore unto him, What sign shewest thou then, that we may see, and believe thee? what dost thou work?"),
        (31, "Our fathers did eat manna in the desert; as it is written, He gave them bread from heaven to eat."),
        (32, "Then Jesus said unto them, Verily, verily, I say unto you, Moses gave you not that bread from heaven; but my Father giveth you the true bread from heaven."),
        (33, "For the bread of God is he which cometh down from heaven, and giveth life unto the world."),
        (34, "Then said they unto him, Lord, evermore give us this bread."),
        (35, "And Jesus said unto them, I am the bread of life: he that cometh to me shall never hunger; and he that believeth on me shall never thirst."),
        (36, "But I said unto you, That ye also have seen me, and believe not."),
        (37, "All that the Father giveth me shall come to me; and him that cometh to me I will in no wise cast out."),
        (38, "For I came down from heaven, not to do mine own will, but the will of him that sent me."),
        (39, "And this is the Father's will which hath sent me, that of all which he hath given me I should lose nothing, but should raise it up again at the last day."),
        (40, "And this is the will of him that sent me, that every one which seeth the Son, and believeth on him, may have everlasting life: and I will raise him up at the last day."),
        (41, "The Jews then murmured at him, because he said, I am the bread which came down from heaven."),
        (42, "And they said, Is not this Jesus, the son of Joseph, whose father and mother we know? how is it then that he saith, I came down from heaven?"),
        (43, "Jesus therefore answered and said unto them, Murmur not among yourselves."),
        (44, "No man can come to me, except the Father which hath sent me draw him: and I will raise him up at the last day."),
        (45, "It is written in the prophets, And they shall be all taught of God. Every man therefore that hath heard, and hath learned of the Father, cometh unto me."),
        (46, "Not that any man hath seen the Father, save he which is of God, he hath seen the Father."),
        (47, "Verily, verily, I say unto you, He that believeth on me hath everlasting life."),
        (48, "I am that bread of life."),
        (49, "Your fathers did eat manna in the wilderness, and are dead."),
        (50, "This is the bread which cometh down from heaven, that a man may eat thereof, and not die."),
        (51, "I am the living bread which came down from heaven: if any man eat of this bread, he shall live for ever: and the bread that I will give is my flesh, which I will give for the life of the world."),
        (52, "The Jews therefore strove among themselves, saying, How can this man give us his flesh to eat?"),
        (53, "Then Jesus said unto them, Verily, verily, I say unto you, Except ye eat the flesh of the Son of man, and drink his blood, ye have no life in you."),
        (54, "Whoso eateth my flesh, and drinketh my blood, hath eternal life; and I will raise him up at the last day."),
        (55, "For my flesh is meat indeed, and my blood is drink indeed."),
        (56, "He that eateth my flesh, and drinketh my blood, dwelleth in me, and I in him."),
        (57, "As the living Father hath sent me, and I live by the Father: so he that eateth me, even he shall live by me."),
        (58, "This is that bread which came down from heaven: not as your fathers did eat manna, and are dead: he that eateth of this bread shall live for ever."),
        (59, "These things said he in the synagogue, as he taught in Capernaum."),
        (60, "Many therefore of his disciples, when they had heard this, said, This is an hard saying; who can hear it?"),
        (61, "When Jesus knew in himself that his disciples murmured at it, he said unto them, Doth this offend you?"),
        (62, "What and if ye shall see the Son of man ascend up where he was before?"),
        (63, "It is the spirit that quickeneth; the flesh profiteth nothing: the words that I speak unto you, they are spirit, and they are life."),
        (64, "But there are some of you that believe not. For Jesus knew from the beginning who they were that believed not, and who should betray him."),
        (65, "And he said, Therefore said I unto you, that no man can come unto me, except it were given unto him of my Father."),
        (66, "From that time many of his disciples went back, and walked no more with him."),
        (67, "Then said Jesus unto the twelve, Will ye also go away?"),
        (68, "Then Simon Peter answered him, Lord, to whom shall we go? thou hast the words of eternal life."),
        (69, "And we believe and are sure that thou art that Christ, the Son of the living God."),
        (70, "Jesus answered them, Have not I chosen you twelve, and one of you is a devil?"),
        (71, "He spake of Judas Iscariot the son of Simon: for he it was that should betray him, being one of the twelve."),
    ]),
]

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)

blank_layout = prs.slide_layouts[6]


def add_section_slide(book_ref):
    slide = prs.slides.add_slide(blank_layout)
    left = Inches(0.5)
    top = Inches(2.5)
    width = prs.slide_width - Inches(1.0)
    height = Inches(2.5)
    tb = slide.shapes.add_textbox(left, top, width, height)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = book_ref
    run.font.size = Pt(54)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)


def add_verse_slide(book_ref, verse_num, verse_text):
    slide = prs.slides.add_slide(blank_layout)

    ref_left = Inches(0.5)
    ref_top = Inches(0.3)
    ref_width = prs.slide_width - Inches(1.0)
    ref_height = Inches(0.8)
    ref_box = slide.shapes.add_textbox(ref_left, ref_top, ref_width, ref_height)
    ref_tf = ref_box.text_frame
    ref_tf.word_wrap = True
    ref_p = ref_tf.paragraphs[0]
    ref_p.alignment = PP_ALIGN.LEFT
    ref_run = ref_p.add_run()
    book = book_ref.split(":")[0].rsplit(" ", 1)[0]
    chapter = book_ref.split(":")[0].rsplit(" ", 1)[1]
    ref_run.text = f"{book} {chapter}:{verse_num}"
    ref_run.font.size = Pt(28)
    ref_run.font.bold = True
    ref_run.font.color.rgb = RGBColor(0x1F, 0x3A, 0x5F)

    body_left = Inches(0.6)
    body_top = Inches(1.3)
    body_width = prs.slide_width - Inches(1.2)
    body_height = prs.slide_height - Inches(1.6)
    body_box = slide.shapes.add_textbox(body_left, body_top, body_width, body_height)
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    body_p = body_tf.paragraphs[0]
    body_p.alignment = PP_ALIGN.LEFT
    num_run = body_p.add_run()
    num_run.text = f"{verse_num} "
    num_run.font.size = Pt(40)
    num_run.font.bold = True
    text_run = body_p.add_run()
    text_run.text = verse_text
    text_run.font.size = Pt(40)


for book_ref, verses in passages:
    add_section_slide(book_ref)
    for v_num, v_text in verses:
        add_verse_slide(book_ref, v_num, v_text)

out_path = r"C:\dev\code\message-search\sermons\breadOfLife\output\breadoflife.pptx"
prs.save(out_path)
print(f"Saved: {out_path}")
print(f"Total slides: {len(prs.slides)}")
